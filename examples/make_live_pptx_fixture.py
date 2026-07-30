"""
실연결 테스트용 .pptx 픽스처를 만드는 스크립트입니다.

tests/test_live_pipeline.py 의 .pptx 분기는 실제 Upstage DocParse 를 타야 하는데,
저장소에는 .pdf 샘플만 있어서 .pptx 를 하나 만들어 둡니다.

    pip install python-pptx
    python examples/make_live_pptx_fixture.py

산출물 fixtures/live_sample.pptx 는 저장소에 커밋합니다.
테스트가 python-pptx 에 의존하지 않게 하려는 의도입니다.
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Pt

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "fixtures" / "live_sample.pptx"

# DocParse 가 텍스트를 그대로 뽑아 왔는지 정확히 확인하기 위한 희귀 문자열.
# 자연스러운 실험 ID 형태로 넣어 LLM 프롬프트를 교란하지 않게 한다.
SENTINEL = "ZEPHYRLOCK-7742"

TITLE = "관성센서 기반 행동인식 연구"
SUBTITLE = "척척발표 실연결 검증용 자료 · 2026"

BODY_SLIDES: list[tuple[str, list[str]]] = [
    (
        "연구 배경",
        [
            "웨어러블 기기의 관성센서(IMU)는 카메라 없이도 사용자 행동을 관측한다.",
            "그러나 IMU 신호는 사람이 읽기 어렵고 라벨을 붙이는 비용이 크다.",
            f"실험 ID: {SENTINEL}",
        ],
    ),
    (
        "제안 방법",
        [
            "영상과 관성센서 신호를 같은 표현 공간으로 사영한다.",
            "대조학습으로 두 모달리티의 짝을 맞추어 라벨 없이 학습한다.",
            "학습된 표현을 행동 분류와 검색에 그대로 재사용한다.",
        ],
    ),
]


def build() -> Path:
    prs = Presentation()

    cover = prs.slides.add_slide(prs.slide_layouts[0])
    cover.shapes.title.text = TITLE
    cover.placeholders[1].text = SUBTITLE

    for title, bullets in BODY_SLIDES:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = title
        frame = slide.placeholders[1].text_frame
        frame.text = bullets[0]
        for line in bullets[1:]:
            para = frame.add_paragraph()
            para.text = line
            para.font.size = Pt(18)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    return OUT


if __name__ == "__main__":
    path = build()
    size_kb = path.stat().st_size / 1024
    print(f"생성: {path.relative_to(ROOT)}  ({size_kb:.1f}KB)")
    print(f"슬라이드 {1 + len(BODY_SLIDES)}장, sentinel={SENTINEL}")
